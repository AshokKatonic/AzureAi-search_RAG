import os
from dotenv import load_dotenv
import logging
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.models import VectorizedQuery
from azure.search.documents.indexes.models import (
    SearchIndex,
    SearchField,
    SearchFieldDataType,
    VectorSearch,
    HnswAlgorithmConfiguration,
    SemanticSearch,
    SemanticConfiguration,
    SemanticPrioritizedFields,
    SemanticField,
)
from openai import AzureOpenAI
logging.getLogger('azure').setLevel(logging.ERROR)

# File processing imports
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation
import openpyxl

# Load environment variables from .env file
load_dotenv()

# --- Configuration ---
# Azure AI Search configuration
AZURE_SEARCH_SERVICE_ENDPOINT = os.getenv("AZURE_SEARCH_SERVICE_ENDPOINT")
AZURE_SEARCH_INDEX_NAME = os.getenv("AZURE_SEARCH_INDEX_NAME")
AZURE_SEARCH_ADMIN_KEY = os.getenv("AZURE_SEARCH_ADMIN_KEY")

# Azure OpenAI configuration
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_CHAT_COMPLETION_DEPLOYED_MODEL_NAME = os.getenv("AZURE_OPENAI_CHAT_COMPLETION_DEPLOYED_MODEL_NAME")
AZURE_OPENAI_EMBEDDING_DEPLOYED_MODEL_NAME = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYED_MODEL_NAME")
AZURE_OPENAI_API_VERSION = "2025-08-07" # Specify a valid API version

# Document chunking configuration
DEFAULT_CHUNK_SIZE = 800  # Optimal size for embeddings (500-1000 chars recommended)
DEFAULT_CHUNK_OVERLAP = 100  # Overlap to maintain context between chunks

# --- Initialize Clients ---
# Initialize Azure OpenAI client
openai_client = AzureOpenAI(
    api_key="3lDUZwn4FGGz4UFp87YH24TgNOaY9M43f0vJHttVORNfWpkv9ziTJQQJ99BIACYeBjFXJ3w3AAABACOGdqM5",
    api_version="2024-02-01",
    azure_endpoint="https://eshal-openai-services.openai.azure.com/"
)

# Initialize Azure AI Search clients
search_credential = AzureKeyCredential(AZURE_SEARCH_ADMIN_KEY)
index_client = SearchIndexClient(endpoint=AZURE_SEARCH_SERVICE_ENDPOINT, credential=search_credential)
search_client = SearchClient(endpoint=AZURE_SEARCH_SERVICE_ENDPOINT, index_name=AZURE_SEARCH_INDEX_NAME, credential=search_credential)


def create_search_index():
    """
    Creates a new search index in Azure AI Search if it doesn't already exist.
    The index is configured for vector search and semantic ranking.
    """
    if AZURE_SEARCH_INDEX_NAME in index_client.list_index_names():
        print(f"Index '{AZURE_SEARCH_INDEX_NAME}' already exists. Skipping creation.")
        return

    print(f"Creating index '{AZURE_SEARCH_INDEX_NAME}'...")
    fields = [
        SearchField(name="id", type=SearchFieldDataType.String, key=True),
        SearchField(name="content", type=SearchFieldDataType.String, searchable=True, retrievable=True),
        SearchField(name="filepath", type=SearchFieldDataType.String, filterable=True, retrievable=True),
        SearchField(
            name="embedding",
            type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
            vector_search_dimensions=1536,  # Dimensions for text-embedding-ada-002
            vector_search_profile_name="my-vector-search-profile",
        ),
    ]

    # Configure HNSW algorithm with custom parameters
    hnsw_config = HnswAlgorithmConfiguration(
        name="my-hnsw-config",
        # Parameters for HNSW algorithm optimization
        parameters={
            "m": 4,  # Number of bi-directional links created for every new element (default: 4)
            "efConstruction": 400,  # Size of dynamic candidate list (default: 400)
            "efSearch": 500,  # Size of dynamic candidate list for search (default: 500)
            "metric": "cosine"  # Distance metric: "cosine", "euclidean", or "dotProduct"
        }
    )
    
    vector_search = VectorSearch(
        profiles=[{"name": "my-vector-search-profile", "algorithm_configuration_name": "my-hnsw-config"}],
        algorithms=[hnsw_config],
    )
    
    semantic_search = SemanticSearch(
        configurations=[
            SemanticConfiguration(
                name="my-semantic-config",
                prioritized_fields=SemanticPrioritizedFields(
                    content_fields=[SemanticField(field_name="content")]
                ),
            )
        ]
    )

    index = SearchIndex(name=AZURE_SEARCH_INDEX_NAME, fields=fields, vector_search=vector_search, semantic_search=semantic_search)
    index_client.create_index(index)
    print(f"Index '{AZURE_SEARCH_INDEX_NAME}' created successfully.")


def generate_embeddings(text):
    """
    Generates embeddings for a given text using Azure OpenAI's embedding model.
    """
    response = openai_client.embeddings.create(
        input=text,
        model=AZURE_OPENAI_EMBEDDING_DEPLOYED_MODEL_NAME
    )
    return response.data[0].embedding


def extract_text_from_pdf(filepath):
    """
    Extracts text content from a PDF file.
    
    Args:
        filepath (str): Path to the PDF file
    
    Returns:
        str: Extracted text content
    """
    try:
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF {filepath}: {e}")
        return ""


def extract_text_from_docx(filepath):
    """
    Extracts text content from a DOCX file.
    
    Args:
        filepath (str): Path to the DOCX file
    
    Returns:
        str: Extracted text content
    """
    try:
        doc = Document(filepath)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from DOCX {filepath}: {e}")
        return ""


def extract_text_from_csv(filepath):
    """
    Extracts text content from a CSV file.
    
    Args:
        filepath (str): Path to the CSV file
    
    Returns:
        str: Extracted text content
    """
    try:
        df = pd.read_csv(filepath)
        # Convert DataFrame to text representation
        text = df.to_string(index=False)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from CSV {filepath}: {e}")
        return ""


def extract_text_from_pptx(filepath):
    """
    Extracts text content from a PPTX file.
    
    Args:
        filepath (str): Path to the PPTX file
    
    Returns:
        str: Extracted text content
    """
    try:
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PPTX {filepath}: {e}")
        return ""


def extract_text_from_xlsx(filepath):
    """
    Extracts text content from an XLSX file.
    
    Args:
        filepath (str): Path to the XLSX file
    
    Returns:
        str: Extracted text content
    """
    try:
        df = pd.read_excel(filepath)
        # Convert DataFrame to text representation
        text = df.to_string(index=False)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from XLSX {filepath}: {e}")
        return ""


def extract_text_from_file(filepath):
    """
    Extracts text content from various file formats.
    
    Args:
        filepath (str): Path to the file
    
    Returns:
        str: Extracted text content
    """
    filename = os.path.basename(filepath)
    file_extension = os.path.splitext(filename)[1].lower()
    
    if file_extension == '.txt':
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading TXT file {filepath}: {e}")
            return ""
    elif file_extension == '.pdf':
        return extract_text_from_pdf(filepath)
    elif file_extension == '.docx':
        return extract_text_from_docx(filepath)
    elif file_extension == '.csv':
        return extract_text_from_csv(filepath)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(filepath)
    elif file_extension in ['.xlsx', '.xls']:
        return extract_text_from_xlsx(filepath)
    else:
        print(f"Unsupported file format: {file_extension}")
        return ""


def chunk_text(text, chunk_size=800, overlap=100):
    """
    Splits text into overlapping chunks of specified size.
    
    Args:
        text (str): The text to chunk
        chunk_size (int): Maximum size of each chunk in characters
        overlap (int): Number of characters to overlap between chunks
    
    Returns:
        list: List of text chunks
    """
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        
        # If this isn't the last chunk, try to break at word boundary
        if end < len(text):
            # Look for the last space before the chunk size limit
            last_space = text.rfind(' ', start, end)
            if last_space > start:  # Found a space within the chunk
                end = last_space
        
        chunk = text[start:end].strip()
        if chunk:  # Only add non-empty chunks
            chunks.append(chunk)
        
        # Move start position with overlap
        start = end - overlap
        if start >= len(text):
            break
    
    return chunks


def process_and_upload_files(directory_path, chunk_size=800, overlap=100):
    """
    Processes all supported files in a given directory, generates embeddings for their
    content, and uploads them to the Azure AI Search index with improved chunking.
    
    Supported formats: TXT, PDF, DOCX, CSV, PPTX, XLSX, XLS
    
    Args:
        directory_path (str): Path to directory containing files
        chunk_size (int): Maximum size of each chunk in characters (default: 800)
        overlap (int): Number of characters to overlap between chunks (default: 100)
    """
    print(f"Processing files in directory: {directory_path}")
    print(f"Using chunk size: {chunk_size} characters, overlap: {overlap} characters")
    
    # Supported file extensions
    supported_extensions = {'.txt', '.pdf', '.docx', '.csv', '.pptx', '.xlsx', '.xls'}
    
    documents = []
    processed_files = 0
    
    for filename in os.listdir(directory_path):
        filepath = os.path.join(directory_path, filename)
        file_extension = os.path.splitext(filename)[1].lower()
        
        # Skip directories and unsupported files
        if os.path.isdir(filepath) or file_extension not in supported_extensions:
            continue
            
        print(f"Processing file: {filename} ({file_extension})")
        
        # Extract text content based on file type
        content = extract_text_from_file(filepath)
        
        if not content.strip():
            print(f"Warning: No content extracted from {filename}")
            continue
            
        # Improved chunking with configurable size and overlap
        chunks = chunk_text(content, chunk_size, overlap)
        print(f"File '{filename}' split into {len(chunks)} chunks")
        
        for i, chunk in enumerate(chunks):
            if chunk.strip(): # Process non-empty chunks
                document = {
                    "id": f"{filename.replace('.', '_')}_{i}",
                    "content": chunk,
                    "filepath": filename,
                    "embedding": generate_embeddings(chunk)
                }
                documents.append(document)
        
        processed_files += 1
    
    print(f"Processed {processed_files} files")
    
    if documents:
        search_client.upload_documents(documents=documents)
        print(f"Uploaded {len(documents)} documents to the index.")
    else:
        print("No new documents to upload.")


def perform_vector_search(query, k=3):
    """
    Performs a vector search on the index for a given query.
    """
    vectorized_query = VectorizedQuery(
    vector=generate_embeddings(query),  # Need to generate embedding first
    k_nearest_neighbors=k,
    fields="embedding"
    )

    results = search_client.search(
        search_text="",  # No text search, only vector search
        vector_queries=[vectorized_query],
        select=["content", "filepath"]
    )
    
    return [result for result in results]


def get_rag_response(query, search_results):
    """
    Generates a response using the RAG pattern. It combines the user's query
    and the search results into a prompt for the chat completion model.
    """
    system_prompt = """
    You are an intelligent assistant. You answer user questions based on the context provided.
    If the information is not in the context, say that you cannot answer.
    """
    
    # Combine search results into a single context string
    context_parts = []
    for result in search_results:
        filepath = result.get('filepath', 'Unknown file')
        content = result.get('content', 'No content available')
        context_parts.append(f"From {filepath}:\n{content}")

    context = "\n\n".join(context_parts)

    # Create the user prompt
    user_prompt = f"Context:\n{context}\n\nQuestion:\n{query}"

    message_text = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt}
    ]

    completion = openai_client.chat.completions.create(
        model=AZURE_OPENAI_CHAT_COMPLETION_DEPLOYED_MODEL_NAME,
        messages=message_text,
        # temperature=0.7,
        max_completion_tokens=800,
    )
    
    return completion.choices[0].message.content

def delete_search_index():
    """
    Deletes the existing search index to allow recreation with proper field configuration.
    
    Returns:
        None
        
    Note:
        - This allows the index to be recreated with retrievable fields
        - Use this when field configuration needs to be updated
    """
    try:
        # First get the index object, then delete it
        if AZURE_SEARCH_INDEX_NAME in index_client.list_index_names():
            index_client.delete_index(index=AZURE_SEARCH_INDEX_NAME)
            print(f"Index '{AZURE_SEARCH_INDEX_NAME}' deleted successfully.")
        else:
            print(f"Index '{AZURE_SEARCH_INDEX_NAME}' does not exist.")
    except Exception as e:
        print(f"Could not delete index: {e}")


if __name__ == "__main__":
    # 1. Create the search index (if it doesn't exist)
    delete_search_index()
    create_search_index()
    
    # 2. Process and upload files from a local directory
    # Create a 'data' folder and add some .txt files to it.
    data_directory = "data" 
    if not os.path.exists(data_directory):
        os.makedirs(data_directory)
        # Create a sample file for demonstration
        with open(os.path.join(data_directory, "sample.txt"), "w") as f:
            f.write("Azure AI Search is a fully managed search-as-a-service. It provides a rich search experience to custom applications.\n\n")
            f.write("Retrieval-Augmented Generation (RAG) is an AI framework for improving the quality of LLM-generated responses. It grounds the model on external sources of knowledge to supplement the LLM’s internal representation of information.")
            
    # Process files with improved chunking
    process_and_upload_files(data_directory, chunk_size=DEFAULT_CHUNK_SIZE, overlap=DEFAULT_CHUNK_OVERLAP)
    
    # 3. Perform a RAG-based query
    print("\n--- Ready to answer questions ---")
    try:
        while True:
            user_query = input("Enter your question (or 'exit' to quit): ")
            if user_query.lower() == 'exit':
                break
            
            # a. Perform vector search
            search_results = perform_vector_search(user_query)
            
            if not search_results:
                print("I couldn't find any relevant information in the documents.")
                continue

            # b. Get RAG response
            response = get_rag_response(user_query, search_results)
            
            print("\nAnswer:")
            print(response)
            print("\n------------------\n")
            
    except KeyboardInterrupt:
        print("\nExiting application.")